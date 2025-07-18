from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from django.conf import settings
from django.core.files.base import ContentFile
import io
import requests
import logging

logger = logging.getLogger(__name__)

class PPTXExportService:
    """Service for exporting presentations to PPTX format"""
    
    def create_pptx(self, presentation_obj):
        """Create PPTX file from presentation object"""
        try:
            # Create presentation
            prs = PPTXPresentation()
            
            # Set slide size (16:9)
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Process each slide
            for slide_obj in presentation_obj.slides.all().order_by('slide_number'):
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title = slide.shapes.title
                title.text = slide_obj.title
                
                # Add content
                content = slide.placeholders[1]
                text_frame = content.text_frame
                text_frame.text = slide_obj.content
                
                # Add image if available
                if slide_obj.image_url:
                    try:
                        # Download image
                        response = requests.get(slide_obj.image_url, timeout=10)
                        if response.status_code == 200:
                            image_stream = io.BytesIO(response.content)
                            
                            # Add image to slide
                            left = Inches(8)
                            top = Inches(2)
                            width = Inches(4)
                            slide.shapes.add_picture(image_stream, left, top, width=width)
                    except Exception as img_error:
                        logger.warning(f"Failed to add image to slide: {str(img_error)}")
            
            # Save to BytesIO
            pptx_stream = io.BytesIO()
            prs.save(pptx_stream)
            pptx_stream.seek(0)
            
            return pptx_stream
            
        except Exception as e:
            logger.error(f"PPTX creation error: {str(e)}")
            raise Exception(f"Failed to create PPTX: {str(e)}")

class PDFExportService:
    """Service for exporting presentations to PDF format"""
    
    def create_pdf(self, presentation_obj):
        """Create PDF file from presentation object"""
        try:
            # Create PDF buffer
            pdf_buffer = io.BytesIO()
            
            # Create document
            doc = SimpleDocTemplate(
                pdf_buffer,
                pagesize=A4,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=18
            )
            
            # Get styles
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                spaceAfter=30,
                alignment=1  # Center alignment
            )
            
            slide_title_style = ParagraphStyle(
                'SlideTitle',
                parent=styles['Heading2'],
                fontSize=18,
                spaceAfter=12,
                textColor='blue'
            )
            
            content_style = ParagraphStyle(
                'SlideContent',
                parent=styles['Normal'],
                fontSize=12,
                spaceAfter=20,
                leftIndent=20
            )
            
            # Build story
            story = []
            
            # Add presentation title
            story.append(Paragraph(presentation_obj.title, title_style))
            story.append(Spacer(1, 12))
            
            if presentation_obj.description:
                story.append(Paragraph(presentation_obj.description, styles['Normal']))
                story.append(Spacer(1, 20))
            
            # Add slides
            for slide_obj in presentation_obj.slides.all().order_by('slide_number'):
                # Add slide title
                story.append(Paragraph(f"Slide {slide_obj.slide_number}: {slide_obj.title}", slide_title_style))
                
                # Add slide content
                story.append(Paragraph(slide_obj.content, content_style))
                
                # Add image if available
                if slide_obj.image_url:
                    try:
                        response = requests.get(slide_obj.image_url, timeout=10)
                        if response.status_code == 200:
                            image_stream = io.BytesIO(response.content)
                            story.append(Image(image_stream, width=4*inch, height=3*inch))
                    except Exception as img_error:
                        logger.warning(f"Failed to add image to PDF: {str(img_error)}")
                
                story.append(Spacer(1, 20))
            
            # Build PDF
            doc.build(story)
            pdf_buffer.seek(0)
            
            return pdf_buffer
            
        except Exception as e:
            logger.error(f"PDF creation error: {str(e)}")
            raise Exception(f"Failed to create PDF: {str(e)}")

# Initialize services
pptx_service = PPTXExportService()
pdf_service = PDFExportService()
